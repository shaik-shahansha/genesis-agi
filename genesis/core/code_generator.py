"""
Intelligent Code Generator - Generates code on-the-fly for any task.

Uses LLM + RAG to create optimized, executable code for user requests.
No pre-built tools - pure dynamic generation!
"""

import ast
import json
import re
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, TYPE_CHECKING

if TYPE_CHECKING:
    from genesis.core.mind import Mind
    from genesis.core.autonomous_orchestrator import UploadedFile


@dataclass
class GeneratedCode:
    """Generated code with metadata."""
    source: str
    language: str
    dependencies: List[str]
    estimated_runtime: float


class IntelligentCodeGenerator:
    """
    Generate optimal code for any task using LLM + RAG.
    
    This is where the magic happens - no pre-built tools,
    just dynamic code generation for ANY request!
    """
    
    def __init__(self, mind: 'Mind'):
        """Initialize code generator."""
        self.mind = mind
    
    async def generate_solution(
        self,
        task: str,
        context: Dict[str, Any],
        files: Optional[List['UploadedFile']] = None
    ) -> GeneratedCode:
        """
        Generate complete, executable code for task.
        
        Args:
            task: Description of what to accomplish
            context: Execution context
            files: Uploaded files to process
            
        Returns:
            GeneratedCode with source and metadata
        """
        
        # Search for similar past solutions
        similar = await self._search_past_solutions(task)
        
        # Build context-aware prompt
        prompt = self._build_prompt(task, context, files, similar)
        
        # Debug: Check what we're working with
        understanding = context.get('understanding', {})
        print(f"[DEBUG CODE_GEN] Task: {task[:80]}")
        print(f"[DEBUG CODE_GEN] Output format: {understanding.get('output_format', 'unknown')}")
        print(f"[DEBUG CODE_GEN] Topic: {understanding.get('topic', 'unknown')}")
        print(f"[DEBUG CODE_GEN] Filename: {understanding.get('output_filename', 'unknown')}")
        
        # Generate code with LLM (use high max_tokens for complete code)
        response = await self.mind.orchestrator.generate(
            messages=[{"role": "user", "content": prompt}],
            model=self.mind.intelligence.reasoning_model,
            max_tokens=4000,  # Ensure we get complete code
            temperature=0.3  # Lower temp for more precise code
        )
        code = response.content
        
        print(f"[CODE GEN] Raw LLM response length: {len(code)}")
        print(f"[CODE GEN] Response preview: {code[:200]}...")
        
        # Extract code from response
        code = self._extract_code_block(code)
        
        print(f"[CODE GEN] Extracted code length: {len(code)}")
        print(f"[CODE GEN] Code preview: {code[:200]}...")
        
        # Validate syntax
        if not self._validate_syntax(code):
            print(f"[CODE GEN]  Syntax validation failed, attempting to fix...")
            # Try to fix automatically
            code = await self._fix_syntax_errors(code)
            
            # Re-validate
            if not self._validate_syntax(code):
                print(f"[CODE GEN]  Still has syntax errors after fix!")
            else:
                print(f"[CODE GEN]  Syntax fixed successfully")
        else:
            print(f"[CODE GEN]  Syntax validated successfully")
        
        return GeneratedCode(
            source=code,
            language="python",
            dependencies=self._extract_dependencies(code),
            estimated_runtime=self._estimate_runtime(code)
        )
    
    def _build_prompt(
        self,
        task: str,
        context: Dict[str, Any],
        files: Optional[List['UploadedFile']],
        similar_solutions: List[Dict]
    ) -> str:
        """Build comprehensive prompt for code generation."""
        
        # Extract context about what to create
        understanding = context.get('understanding', {})
        output_format = understanding.get('output_format', 'other')
        topic = understanding.get('topic', task)
        filename = understanding.get('output_filename', 'output')
        
        # Fallback: If no understanding context, try to infer from task
        if not understanding or output_format == 'other':
            import re
            task_lower = task.lower()
            if any(word in task_lower for word in ['document', 'doc', 'word']):
                output_format = 'document'
            elif any(word in task_lower for word in ['presentation', 'ppt', 'powerpoint', 'slides']):
                output_format = 'presentation'
            elif any(word in task_lower for word in ['report', 'analysis report']):
                output_format = 'report'
            
            # Generate filename from task if not provided
            if filename == 'output':
                # Extract topic from task
                topic_match = re.search(r'(?:about|on|for)\s+(.+?)(?:\s+with|$)', task_lower)
                if topic_match:
                    topic_text = topic_match.group(1).strip()
                    filename = re.sub(r'[^\w\s-]', '', topic_text)
                    filename = re.sub(r'[-\s]+', '_', filename)
                    filename = filename[:50]
        
        prompt = f"""Generate Python code to accomplish this task:
{task}

Context:
{json.dumps(context, indent=2)}

"""
        
        # Add file information if present
        if files:
            prompt += "\nFiles available:\n"
            for file in files:
                prompt += f"- {file.name} ({file.mime_type}, {file.size} bytes)\n"
                prompt += f"  Path: {file.path}\n"
        
        # Add format-specific guidance
        if output_format == "document":
            prompt += f"""\n\nDOCUMENT GENERATION REQUIREMENTS:
- Create a professional Word document (.docx) about: {topic}
- Filename: {filename}.docx
- Use python-docx library
- Include proper structure: title, sections, paragraphs
- Add professional formatting: headings, fonts, spacing
- Content should be comprehensive and relevant to the topic
- DO NOT create random content - focus on {topic}

Template structure:
1. Title page with document title
2. Introduction section
3. Main content sections (2-4 sections relevant to topic)
4. Conclusion
5. Professional formatting throughout
"""
        elif output_format == "presentation":
            prompt += f"""\n\nPRESENTATION GENERATION REQUIREMENTS:
- Create a professional PowerPoint presentation (.pptx) about: {topic}
- Filename: {filename}.pptx
- Use python-pptx library
- Include 5-8 slides with proper structure
- Professional design with consistent formatting
- Content should be focused on {topic}
- DO NOT create random content

Template structure:
1. Title slide with topic name
2. Overview/Agenda slide
3. Content slides (3-5 slides) with key points
4. Each slide should have:
   - Clear title
   - Bullet points or text content (2-5 key points per slide)
   - Professional layout
5. Conclusion/Summary slide

Layout guidelines:
- Use built-in layouts: presentation.slide_layouts[0] for title, [1] for title+content, [6] for blank
- Add text boxes using slide.shapes.title and slide.placeholders
- Use proper text formatting with text_frame.text and paragraphs
- Keep it simple and professional

Example code pattern:
```python
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

# Title slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Your Topic"
subtitle = title_slide.placeholders[1]
subtitle.text = "Subtitle text"

# Content slide
content_slide = prs.slides.add_slide(prs.slide_layouts[1])
content_slide.shapes.title.text = "Slide Title"
content = content_slide.placeholders[1].text_frame
content.text = "First point"
for point in ["Second point", "Third point"]:
    p = content.add_paragraph()
    p.text = point
    p.level = 0

prs.save('{filename}.pptx')
```

IMPORTANT: DO NOT try to add images - they cause errors. Focus on text content and layout only.
"""
        elif output_format == "report":
            prompt += f"""\n\nREPORT GENERATION REQUIREMENTS:
- Create a professional report document (.docx) about: {topic}
- Filename: {filename}.docx
- Use python-docx library
- Include executive summary, detailed sections, and data/findings
- Professional business report formatting
- Content focused on {topic}
"""
        
        prompt += """
Available libraries (will be auto-installed if needed):
- requests (for API calls, image downloads)
- python-pptx (PowerPoint generation)
- python-docx (Word document generation)
- openpyxl (Excel files)
- pandas, numpy (data analysis)
- pillow (image processing)
- matplotlib, seaborn (charts/graphs)
- urllib.parse (URL encoding)

Code structure:
```python
import json
import sys
from pathlib import Path
# Other imports as needed...

def main():
    \"\"\"Main function to accomplish the task.\"\"\"
    try:
        # CRITICAL: Use the exact filename specified above
        # For this task, the filename is: {filename}
        # Create the file with the correct name and extension
        
        # Your code here to generate the artifact
        # Save to current directory
        # Print success with file path
        
        result = {{
            "success": True,
            "message": "Task completed successfully",
            "data": {{}},  # Include relevant outputs
            "generated_files": []  # List of file paths created
        }}
        print(json.dumps(result))
        
    except Exception as e:
        print(json.dumps({{
            "success": False,
            "error": str(e),
            "error_type": type(e).__name__
        }}), file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
```

CRITICAL INSTRUCTIONS:
1. Return ONLY executable Python code - no explanations, no markdown except code block
2. DO NOT return JSON metadata or thinking steps
3. Create EXACTLY what the user asked for - topic is: {topic}
4. Use the EXACT filename: {filename} (with appropriate extension)
5. DO NOT use generic names like 'output.docx' or 'generated_document.docx'
6. For presentations: Create text-based slides with clear structure and bullet points
7. For documents: MUST include proper structure and formatting
8. Save files to current directory (use Path.cwd())
9. Print clear success message with file location
10. Handle all errors gracefully
11. CRITICAL: DO NOT try to download or add images to presentations - focus on content only

NOTE: If filename already exists, system will automatically version it (e.g., {filename}_v2.docx)
"""
        
        return prompt
    
    async def _search_past_solutions(self, task: str) -> List[Dict]:
        """Search for similar past solutions in memory."""
        try:
            results = await self.mind.memory.search(
                query=task,
                memory_type="procedural",
                k=3
            )
            return results
        except Exception:
            return []
    
    def _extract_code_block(self, response: str) -> str:
        """Extract code from markdown code blocks or JSON-wrapped responses."""
        
        # First, check if response is wrapped in JSON (common LLM mistake)
        try:
            # Try to parse as JSON
            data = json.loads(response)
            
            # If it's a list with "type" fields, extract code from appropriate entry
            if isinstance(data, list):
                for item in data:
                    if isinstance(item, dict) and item.get('type') in ['code', 'execute', 'python']:
                        if 'code' in item:
                            return item['code']
                        elif 'content' in item:
                            return item['content']
            
            # If it's a dict with code field
            if isinstance(data, dict) and 'code' in data:
                return data['code']
                
        except (json.JSONDecodeError, TypeError):
            pass
        
        # Try to find Python code block
        patterns = [
            r"```python\n(.*?)\n```",  # python block
            r"```py\n(.*?)\n```",      # py block  
            r"```\n(.*?)\n```",        # generic block
        ]
        
        for pattern in patterns:
            match = re.search(pattern, response, re.DOTALL)
            if match:
                code = match.group(1).strip()
                # Verify it looks like Python
                if any(keyword in code for keyword in ['import ', 'def ', 'class ', 'if ', 'for ']):
                    return code
        
        # If no code block found, try to extract Python-looking content
        lines = response.split('\n')
        code_lines = []
        in_code = False
        
        for line in lines:
            # Start collecting when we see Python keywords
            if not in_code and any(line.strip().startswith(kw) for kw in ['import ', 'from ', 'def ', 'class ']):
                in_code = True
            
            if in_code:
                # Stop if we hit markdown or JSON
                if line.strip().startswith(('```', '{', '[')):
                    if not line.strip().startswith('```python'):
                        break
                code_lines.append(line)
        
        if code_lines:
            return '\n'.join(code_lines).strip()
        
        # Last resort: return entire response stripped
        return response.strip()
    
    def _validate_syntax(self, code: str) -> bool:
        """Validate Python syntax."""
        try:
            ast.parse(code)
            return True
        except SyntaxError:
            return False
    
    async def _fix_syntax_errors(self, code: str) -> str:
        """Try to fix syntax errors automatically."""
        prompt = f"""This Python code has syntax errors. Fix them:

```python
{code}
```

Return ONLY the corrected code, no explanations."""
        
        response = await self.mind.orchestrator.generate(
            messages=[{"role": "user", "content": prompt}],
            model=self.mind.intelligence.reasoning_model,
            max_tokens=4000,
            temperature=0.3
        )
        return self._extract_code_block(response.content)
    
    def _extract_dependencies(self, code: str) -> List[str]:
        """Extract import statements from code."""
        dependencies = []
        
        # Parse imports
        try:
            tree = ast.parse(code)
            for node in ast.walk(tree):
                if isinstance(node, ast.Import):
                    for alias in node.names:
                        dependencies.append(alias.name)
                elif isinstance(node, ast.ImportFrom):
                    if node.module:
                        dependencies.append(node.module)
        except Exception:
            pass
        
        return list(set(dependencies))
    
    def _estimate_runtime(self, code: str) -> float:
        """Estimate code execution time (rough heuristic)."""
        # Simple heuristic based on code complexity
        lines = len(code.split('\n'))
        
        # Look for expensive operations
        if 'requests.get' in code or 'BeautifulSoup' in code:
            return 10.0  # Web scraping
        elif 'pandas' in code and lines > 50:
            return 5.0  # Data processing
        elif 'playwright' in code:
            return 15.0  # Browser automation
        else:
            return 2.0  # Simple task
